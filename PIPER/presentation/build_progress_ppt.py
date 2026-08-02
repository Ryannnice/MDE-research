#!/usr/bin/env python3
"""Build a concise, image-backed PPTX for the PIPER project progress."""

from __future__ import annotations

from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches


ROOT = Path(__file__).resolve().parents[2]
OUT_DIR = Path(__file__).resolve().parent
SLIDE_DIR = OUT_DIR / "slides"
PPTX_PATH = OUT_DIR / "PIPER_全流程阶段进展_2026-08-02.pptx"
PREVIEW_PATH = OUT_DIR / "PIPER_全流程阶段进展_预览.png"
FONT_PATH = str(next(path for path in (
    Path("/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc"),
    Path("/tmp/piper_noto_cjk/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc"),
    Path("/usr/share/fonts/truetype/droid/DroidSansFallbackFull.ttf"),
) if path.exists()))

W, H = 1920, 1080
BG = "#0B1220"
CARD = "#121D31"
CARD_2 = "#17243B"
TEXT = "#F4F7FB"
MUTED = "#AAB7CC"
CYAN = "#22D3EE"
BLUE = "#60A5FA"
GREEN = "#34D399"
AMBER = "#FBBF24"
RED = "#F87171"
LINE = "#2A3B57"


def font(size: int) -> ImageFont.FreeTypeFont:
    return ImageFont.truetype(FONT_PATH, size=size)


def canvas() -> tuple[Image.Image, ImageDraw.ImageDraw]:
    image = Image.new("RGB", (W, H), BG)
    draw = ImageDraw.Draw(image)
    draw.ellipse((1450, -320, 2180, 410), fill="#102A43")
    draw.ellipse((-350, 820, 350, 1520), fill="#101D33")
    return image, draw


def card(draw: ImageDraw.ImageDraw, box: tuple[int, int, int, int], fill: str = CARD,
         outline: str | None = None, radius: int = 28) -> None:
    draw.rounded_rectangle(box, radius=radius, fill=fill, outline=outline, width=2 if outline else 1)


def pill(draw: ImageDraw.ImageDraw, x: int, y: int, text: str, color: str,
         text_color: str = BG, size: int = 25) -> int:
    f = font(size)
    bbox = draw.textbbox((0, 0), text, font=f)
    width = bbox[2] - bbox[0] + 44
    draw.rounded_rectangle((x, y, x + width, y + 48), radius=24, fill=color)
    draw.text((x + 22, y + 22), text, font=f, fill=text_color, anchor="lm")
    return width


def title(draw: ImageDraw.ImageDraw, index: str, heading: str, subtitle: str = "") -> None:
    draw.text((80, 78), index, font=font(28), fill=CYAN)
    draw.text((80, 120), heading, font=font(58), fill=TEXT)
    if subtitle:
        draw.text((82, 195), subtitle, font=font(27), fill=MUTED)
    draw.line((80, 250, 1840, 250), fill=LINE, width=2)


def footer(draw: ImageDraw.ImageDraw, page: int) -> None:
    draw.text((80, 1030), "PIPER ROBOT PLATFORM  ·  2026.08.02", font=font(20), fill="#6F809A")
    draw.text((1840, 1030), f"{page:02d}", font=font(22), fill="#6F809A", anchor="ra")


def arrow(draw: ImageDraw.ImageDraw, start: tuple[int, int], end: tuple[int, int],
          color: str = CYAN, width: int = 5) -> None:
    draw.line((start, end), fill=color, width=width)
    x2, y2 = end
    x1, y1 = start
    if abs(x2 - x1) >= abs(y2 - y1):
        direction = 1 if x2 > x1 else -1
        head = [(x2, y2), (x2 - direction * 18, y2 - 11), (x2 - direction * 18, y2 + 11)]
    else:
        direction = 1 if y2 > y1 else -1
        head = [(x2, y2), (x2 - 11, y2 - direction * 18), (x2 + 11, y2 - direction * 18)]
    draw.polygon(head, fill=color)


def multiline(draw: ImageDraw.ImageDraw, xy: tuple[int, int], lines: list[str],
              size: int = 29, color: str = TEXT, spacing: int = 16,
              bullet_color: str = GREEN) -> None:
    x, y = xy
    f = font(size)
    for line in lines:
        draw.ellipse((x, y + 10, x + 10, y + 20), fill=bullet_color)
        draw.text((x + 28, y), line, font=f, fill=color)
        y += size + spacing


def fit_image(source: Path, size: tuple[int, int]) -> Image.Image:
    image = Image.open(source).convert("RGB")
    target_w, target_h = size
    scale = max(target_w / image.width, target_h / image.height)
    resized = image.resize((round(image.width * scale), round(image.height * scale)), Image.Resampling.LANCZOS)
    left = (resized.width - target_w) // 2
    top = (resized.height - target_h) // 2
    return resized.crop((left, top, left + target_w, top + target_h))


def slide_cover() -> Image.Image:
    image, draw = canvas()
    draw.rounded_rectangle((80, 92, 262, 143), radius=25, fill=CYAN)
    draw.text((171, 118), "阶段进展", font=font(25), fill=BG, anchor="mm")
    draw.text((80, 245), "从铝制实验桌到", font=font(73), fill=TEXT)
    draw.text((80, 340), "远程机器人平台", font=font(88), fill=CYAN)
    draw.text((84, 470), "松灵 PIPER · AGILE.X 夹爪 · Intel RealSense D455", font=font(34), fill=MUTED)

    steps = [
        ("01", "实验桌"), ("02", "机械臂"), ("03", "末端设备"),
        ("04", "远程连接"), ("05", "RGB-D 感知"),
    ]
    x = 82
    for idx, (number, label) in enumerate(steps):
        draw.ellipse((x, 680, x + 66, 746), fill=CYAN if idx < 5 else LINE)
        draw.text((x + 33, 713), number, font=font(23), fill=BG, anchor="mm")
        draw.text((x + 33, 770), label, font=font(28), fill=TEXT, anchor="ma")
        if idx < len(steps) - 1:
            draw.line((x + 76, 713, x + 260, 713), fill=LINE, width=4)
        x += 330

    card(draw, (80, 882, 1840, 977), fill="#10263A", outline="#1D5262", radius=25)
    draw.text((120, 930), "当前结论", font=font(26), fill=CYAN, anchor="lm")
    draw.text((320, 930), "硬件平台、机械臂、夹爪、远程链路和 D455 已全部打通", font=font(34), fill=TEXT, anchor="lm")
    return image


def draw_robot_on_table(draw: ImageDraw.ImageDraw) -> None:
    # Aluminium workbench.
    draw.rounded_rectangle((120, 555, 880, 625), radius=16, fill="#7D8DA5", outline="#BDD0E8", width=3)
    for x in (165, 790):
        draw.polygon([(x, 625), (x + 55, 625), (x + 25, 930), (x - 20, 930)], fill="#506179")
    draw.line((190, 780, 805, 780), fill="#506179", width=32)
    draw.line((205, 650, 780, 900), fill="#405068", width=14)
    draw.line((780, 650, 210, 900), fill="#405068", width=14)

    # Simplified PIPER arm.
    draw.rounded_rectangle((440, 500, 600, 570), radius=20, fill="#D6DEE8")
    joints = [(520, 500), (535, 395), (655, 315), (765, 365), (835, 290)]
    for a, b in zip(joints, joints[1:]):
        draw.line((a, b), fill="#E9EEF5", width=34)
        draw.line((a, b), fill="#8FA3BA", width=6)
    for x, y in joints:
        draw.ellipse((x - 26, y - 26, x + 26, y + 26), fill="#C7D2DF", outline=CYAN, width=4)
    draw.line((835, 290, 885, 270), fill="#D8E0E9", width=22)
    draw.line((883, 270, 912, 245), fill="#D8E0E9", width=12)
    draw.line((883, 270, 917, 288), fill="#D8E0E9", width=12)
    draw.rounded_rectangle((760, 240, 840, 276), radius=8, fill="#263D57", outline=CYAN, width=3)


def slide_physical() -> Image.Image:
    image, draw = canvas()
    title(draw, "01", "实体平台搭建", "先完成稳定工作台，再逐步安装机器人与末端设备")
    card(draw, (80, 290, 970, 970), fill="#101B2D", outline=LINE)
    draw_robot_on_table(draw)
    draw.text((525, 920), "铝制实验桌 + PIPER + 夹爪 + 腕部相机", font=font(27), fill=MUTED, anchor="mm")

    items = [
        ("1", "铝制实验桌", "完成框架安装、桌面固定与整体紧固"),
        ("2", "PIPER 上桌", "机械臂底座固定，接通电源并确认安装稳定"),
        ("3", "AGILE.X 夹爪", "安装末端执行器，接入机械臂 24 V / CAN 线束"),
        ("4", "D455 相机", "固定在末端附近，USB 3.x 线直接连接 Windows"),
    ]
    y = 305
    colors = [CYAN, BLUE, GREEN, AMBER]
    for (num, name, desc), color in zip(items, colors):
        card(draw, (1040, y, 1840, y + 145), fill=CARD, outline=LINE, radius=22)
        draw.ellipse((1080, y + 38, 1150, y + 108), fill=color)
        draw.text((1115, y + 73), num, font=font(28), fill=BG, anchor="mm")
        draw.text((1190, y + 35), name, font=font(34), fill=TEXT)
        draw.text((1192, y + 87), desc, font=font(25), fill=MUTED)
        y += 165
    footer(draw, 2)
    return image


def system_box(draw: ImageDraw.ImageDraw, box: tuple[int, int, int, int],
               label: str, detail: str, accent: str) -> None:
    card(draw, box, fill=CARD, outline=accent, radius=22)
    x1, y1, x2, y2 = box
    draw.rectangle((x1, y1, x1 + 8, y2), fill=accent)
    draw.text(((x1 + x2) // 2, y1 + 47), label, font=font(31), fill=TEXT, anchor="mm")
    draw.text(((x1 + x2) // 2, y1 + 95), detail, font=font(22), fill=MUTED, anchor="mm")


def slide_architecture() -> Image.Image:
    image, draw = canvas()
    title(draw, "02", "完整连接链路", "本地 VS Code 通过学校服务器进入现场 Windows，再连接真实硬件")

    boxes = [
        ((70, 395, 330, 535), "个人电脑", "VS Code / SSH", BLUE),
        ((405, 395, 710, 535), "学校服务器", "3090x6", CYAN),
        ((840, 365, 1190, 565), "现场 Windows", "ENVY_Katana", GREEN),
        ((1300, 395, 1515, 535), "USB-CAN", "1 Mbps", AMBER),
        ((1610, 395, 1850, 535), "PIPER", "S-V1.8-9", CYAN),
    ]
    for box, label, detail, accent in boxes:
        system_box(draw, box, label, detail, accent)
    arrow(draw, (330, 465), (405, 465), BLUE)
    arrow(draw, (710, 465), (840, 465), CYAN)
    arrow(draw, (1190, 465), (1300, 465), GREEN)
    arrow(draw, (1515, 465), (1610, 465), AMBER)

    draw.text((775, 410), "反向 SSH", font=font(21), fill=CYAN, anchor="ma")
    draw.text((775, 435), "127.0.0.1:22022", font=font(18), fill=MUTED, anchor="ma")

    system_box(draw, (845, 700, 1175, 850), "D455", "USB 3.2 · RGB-D/IMU", BLUE)
    arrow(draw, (1010, 700), (1010, 575), BLUE)
    system_box(draw, (1585, 700, 1860, 850), "AGILE.X 夹爪", "24 V / CAN", GREEN)
    arrow(draw, (1722, 700), (1722, 545), GREEN)

    card(draw, (75, 705, 710, 900), fill="#10263A", outline="#1D5262")
    draw.text((115, 750), "Windows 上的两个常驻任务", font=font(31), fill=CYAN)
    multiline(draw, (120, 810), [
        "PIPER Bridge Observe：127.0.0.1:57845",
        "PIPER Remote Tunnel：断线自动重连",
    ], size=25, spacing=16, bullet_color=GREEN)
    footer(draw, 3)
    return image


def slide_arm_gripper() -> Image.Image:
    image, draw = canvas()
    title(draw, "03", "机械臂与夹爪联调", "从“读到状态”推进到“可控、可验证、可恢复”")

    card(draw, (80, 300, 920, 850), fill=CARD, outline=CYAN)
    draw.text((125, 350), "松灵 PIPER", font=font(42), fill=CYAN)
    pill(draw, 660, 340, "已验证", GREEN)
    multiline(draw, (130, 435), [
        "USB-CAN 通信打通，固件 S-V1.8-9",
        "六关节反馈、状态与错误码稳定读取",
        "完成 J2/J3、J6 小角度分级动作",
        "完成“竖直指向上方—归位”循环",
        "当前保持 observe，只读且无活动命令",
    ], size=29, spacing=22, bullet_color=CYAN)

    card(draw, (1000, 300, 1840, 850), fill=CARD, outline=GREEN)
    draw.text((1045, 350), "AGILE.X 原厂夹爪", font=font(42), fill=GREEN)
    pill(draw, 1580, 340, "已验证", GREEN)
    multiline(draw, (1050, 435), [
        "反馈链路与所有故障位检查通过",
        "完成手动闭合、写零与小行程张开",
        "支持 0–70 mm 宽度、0–3 N 力参数",
        "动作后关闭驱动，当前开口约 9.8 mm",
        "已兼容该固件 homed 位反馈特点",
    ], size=29, spacing=22, bullet_color=GREEN)

    metrics = [
        ("6 / 6", "关节在线", CYAN),
        ("0", "机械臂错误码", GREEN),
        ("37", "自动测试通过", BLUE),
        ("0", "当前活动命令", AMBER),
    ]
    x = 80
    for value, label, color in metrics:
        card(draw, (x, 885, x + 410, 990), fill=CARD_2, radius=20)
        draw.text((x + 35, 936), value, font=font(43), fill=color, anchor="lm")
        draw.text((x + 180, 936), label, font=font(25), fill=MUTED, anchor="lm")
        x += 440
    footer(draw, 4)
    return image


def slide_remote_pitfalls() -> Image.Image:
    image, draw = canvas()
    title(draw, "04", "远程链路与踩坑", "最终结果：不再依赖手工保持 PowerShell 窗口")

    phases = [
        ("依赖安装", "GitHub 访问失败", "固定源码并离线打包", BLUE),
        ("CAN 接入", "设备打开卡住", "排除旧进程与 USB 句柄", CYAN),
        ("HTTP 桥", "端口 10013 / 401", "修正端口；保留令牌保护", AMBER),
        ("Windows SSH", "在线安装慢、服务 1067", "部署离线 Win32-OpenSSH", GREEN),
        ("反向隧道", "22022 陈旧监听", "服务器工具释放并自动重连", RED),
    ]
    y = 300
    for idx, (phase, problem, solution, color) in enumerate(phases, start=1):
        draw.ellipse((102, y + 17, 168, y + 83), fill=color)
        draw.text((135, y + 50), str(idx), font=font(27), fill=BG, anchor="mm")
        if idx < len(phases):
            draw.line((135, y + 83, 135, y + 142), fill=LINE, width=5)
        card(draw, (205, y, 1815, y + 100), fill=CARD, outline=LINE, radius=20)
        draw.text((245, y + 20), phase, font=font(28), fill=color)
        draw.text((535, y + 22), problem, font=font(27), fill=TEXT)
        arrow(draw, (930, y + 50), (1020, y + 50), color, width=4)
        draw.text((1070, y + 22), solution, font=font(27), fill=TEXT)
        y += 125

    card(draw, (205, 930, 1815, 990), fill="#10263A", outline="#1D5262", radius=20)
    draw.text((1010, 960), "当前：密钥登录 + 开机计划任务 + 断线恢复工具，服务器可直接 ssh piper-windows", font=font(27), fill=CYAN, anchor="mm")
    footer(draw, 5)
    return image


def crop_depth(source: Path) -> Image.Image:
    image = Image.open(source).convert("RGB")
    return image.crop((image.width // 2, 0, image.width, image.height))


def slide_camera() -> Image.Image:
    image, draw = canvas()
    title(draw, "05", "D455 深度相机联调", "末端腕部视角：RGB、对齐深度与 IMU 已同时采集")

    color_path = ROOT / "PIPER/camera/captures/baseline_v2/color.png"
    color = fit_image(color_path, (820, 465))
    image.paste(color, (80, 300))
    draw.rounded_rectangle((80, 300, 900, 765), radius=25, outline=CYAN, width=4)
    draw.text((110, 725), "真实腕部视角：画面下缘可见夹爪", font=font(23), fill=TEXT)

    info = [
        ("设备", "Intel RealSense D455"),
        ("连接", "USB 3.2 · Firmware 5.17.3.10"),
        ("RGB / 深度", "848 × 480 @ 30 Hz"),
        ("IMU", "加速度计 + 陀螺仪 @ 200 Hz"),
        ("同步", "RGB/深度时间差 < 0.6 ms"),
    ]
    y = 305
    for key, value in info:
        card(draw, (980, y, 1840, y + 78), fill=CARD, radius=18)
        draw.text((1020, y + 39), key, font=font(25), fill=CYAN, anchor="lm")
        draw.text((1260, y + 39), value, font=font(26), fill=TEXT, anchor="lm")
        y += 92

    labels = ["Custom", "High Density", "Shift 50", "Shift 100"]
    values = [24.33, 22.40, 28.51, 40.78]
    colors = [BLUE, "#7C8CA4", CYAN, GREEN]
    draw.text((980, 795), "当前画面有效深度像素", font=font(28), fill=TEXT)
    x0, y0, max_w = 1205, 847, 555
    for label, value, color in zip(labels, values, colors):
        draw.text((980, y0), label, font=font(22), fill=MUTED, anchor="lm")
        draw.rounded_rectangle((x0, y0 - 12, x0 + max_w, y0 + 12), radius=12, fill="#27354B")
        draw.rounded_rectangle((x0, y0 - 12, x0 + int(max_w * value / 50), y0 + 12), radius=12, fill=color)
        draw.text((1795, y0), f"{value:.2f}%", font=font(22), fill=color, anchor="rm")
        y0 += 45

    card(draw, (80, 805, 900, 975), fill="#10263A", outline="#1D5262", radius=22)
    draw.text((120, 845), "使用建议", font=font(29), fill=CYAN)
    draw.text((120, 900), "接近阶段：Custom / shift 50", font=font(25), fill=TEXT)
    draw.text((485, 900), "最后近抓取：shift 100（约 0.18–0.41 m）", font=font(25), fill=GREEN)
    draw.text((120, 944), "所有测试结束后均自动恢复相机原始参数", font=font(23), fill=MUTED)
    footer(draw, 6)
    return image


def slide_status_next() -> Image.Image:
    image, draw = canvas()
    title(draw, "06", "当前残留问题与下一步", "核心功能已打通，但实验环境和本地硬件仍需补齐")

    issues = [
        ("1", "实验室没有网络", "软件下载、远程运维与数据回传受限，\n需要补充稳定网络。", RED),
        ("2", "实验桌晃动明显", "机械臂运动会带动末端相机抖动，\n影响图像、深度与标定精度。", RED),
        ("3", "缺少 USB 扩展坞", "USB-CAN、D455 与第二相机并行接入时，\n电脑端口不足。", AMBER),
        ("4", "缺少本地主机", "目前采用“学校服务器 + 现场 Windows”\n的远程过渡方案。", AMBER),
    ]
    positions = [
        (80, 305, 920, 535),
        (1000, 305, 1840, 535),
        (80, 570, 920, 800),
        (1000, 570, 1840, 800),
    ]
    for (num, heading, detail, color), box in zip(issues, positions):
        card(draw, box, fill=CARD, outline=color, radius=24)
        x1, y1, _, _ = box
        draw.ellipse((x1 + 38, y1 + 38, x1 + 108, y1 + 108), fill=color)
        draw.text((x1 + 73, y1 + 73), num, font=font(28), fill=BG, anchor="mm")
        draw.text((x1 + 140, y1 + 40), heading, font=font(37), fill=color)
        draw.multiline_text((x1 + 42, y1 + 130), detail, font=font(25), fill=TEXT, spacing=8)

    card(draw, (80, 845, 1840, 995), fill="#10263A", outline=CYAN, radius=24)
    draw.text((125, 885), "下一步", font=font(28), fill=CYAN)
    draw.text((345, 885), "手眼标定（camera ↔ gripper）", font=font(42), fill=TEXT)
    draw.text((125, 947), "前提：先加固实验桌和相机安装，保证采集过程中画面与外参稳定。", font=font(27), fill=MUTED)
    footer(draw, 7)
    return image


def build() -> None:
    SLIDE_DIR.mkdir(parents=True, exist_ok=True)
    slides = [
        slide_cover(),
        slide_physical(),
        slide_architecture(),
        slide_arm_gripper(),
        slide_remote_pitfalls(),
        slide_camera(),
        slide_status_next(),
    ]

    slide_paths: list[Path] = []
    for index, slide in enumerate(slides, start=1):
        path = SLIDE_DIR / f"slide_{index:02d}.png"
        slide.save(path, quality=95)
        slide_paths.append(path)

    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    blank = presentation.slide_layouts[6]
    for path in slide_paths:
        slide = presentation.slides.add_slide(blank)
        slide.shapes.add_picture(str(path), 0, 0, width=presentation.slide_width, height=presentation.slide_height)
    presentation.save(PPTX_PATH)

    thumb_w, thumb_h = 640, 360
    preview = Image.new("RGB", (thumb_w * 2, thumb_h * 4), "#050A12")
    for index, slide in enumerate(slides):
        thumb = slide.resize((thumb_w, thumb_h), Image.Resampling.LANCZOS)
        preview.paste(thumb, ((index % 2) * thumb_w, (index // 2) * thumb_h))
    preview.save(PREVIEW_PATH)

    print(PPTX_PATH)
    print(PREVIEW_PATH)


if __name__ == "__main__":
    build()
