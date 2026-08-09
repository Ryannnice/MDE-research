#!/usr/bin/env python3
"""Generate the concise nine-slide transparent-depth project deck."""

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


HERE = Path(__file__).resolve().parent
ASSET_DIR = HERE / "assets"
OUTPUT = HERE / "透明物体多层深度_简洁汇报_2026-08-08.pptx"

SLIDE_W = 13.333
SLIDE_H = 7.5
FONT = "Microsoft YaHei"

INK = "172331"
MUTED = "667381"
TEAL = "0F766E"
BLUE = "2563EB"
GREEN = "15805D"
ORANGE = "C46B16"
RED = "B7443E"
LINE = "DDE3E8"
SOFT = "F5F8FA"
SOFT_TEAL = "EAF6F4"
SOFT_BLUE = "EEF4FF"
SOFT_ORANGE = "FFF4E8"
WHITE = "FFFFFF"


def color(value):
    return RGBColor.from_string(value)


def crop_image(source, target, relative_box):
    source_path = ASSET_DIR / source
    target_path = ASSET_DIR / target
    image = Image.open(source_path)
    width, height = image.size
    left, top, right, bottom = relative_box
    image.crop((
        int(left * width), int(top * height),
        int(right * width), int(bottom * height),
    )).save(target_path)
    return target_path


def textbox(slide, text, x, y, w, h, size=14, fill=INK, bold=False,
            align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, margin=0.02):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color(fill)
    return shape


def rect(slide, x, y, w, h, fill=WHITE, outline=LINE, rounded=False, width=0.8):
    kind = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color(fill)
    shape.line.color.rgb = color(outline)
    shape.line.width = Pt(width)
    return shape


def line(slide, x1, y1, x2, y2, stroke=LINE, width=1.0):
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    connector.line.color.rgb = color(stroke)
    connector.line.width = Pt(width)
    return connector


def circle(slide, x, y, diameter, fill, label=None, label_color=WHITE):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(diameter), Inches(diameter)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color(fill)
    shape.line.color.rgb = color(fill)
    if label is not None:
        textbox(slide, label, x, y + diameter * 0.29, diameter, diameter * 0.3,
                10, label_color, True, PP_ALIGN.CENTER)
    return shape


def picture_contain(slide, path, x, y, w, h):
    with Image.open(path) as image:
        image_w, image_h = image.size
    scale = min(w / image_w, h / image_h)
    draw_w = image_w * scale
    draw_h = image_h * scale
    draw_x = x + (w - draw_w) / 2
    draw_y = y + (h - draw_h) / 2
    slide.shapes.add_picture(
        str(path), Inches(draw_x), Inches(draw_y), Inches(draw_w), Inches(draw_h)
    )
    border = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    border.fill.background()
    border.line.color.rgb = color(LINE)
    border.line.width = Pt(0.7)


def background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color(WHITE)


def title(slide, text, index):
    textbox(slide, f"0{index}", 0.56, 0.37, 0.50, 0.20, 9, TEAL, True)
    textbox(slide, text, 1.05, 0.31, 11.4, 0.42, 24, INK, True)
    line(slide, 0.56, 0.93, 12.74, 0.93, LINE, 0.9)


def footer(slide, text):
    textbox(slide, text, 0.57, 7.18, 11.9, 0.15, 7.2, MUTED)


def chip(slide, text, x, y, w, fill_color, text_color):
    rect(slide, x, y, w, 0.32, fill_color, fill_color, rounded=True)
    textbox(slide, text, x, y + 0.075, w, 0.13, 8.8, text_color, True, PP_ALIGN.CENTER)


def build():
    # Paper figures cropped from the local PDFs.
    transcg = crop_image("transcg_p1.png", "transcg_core.png", (0.50, 0.24, 0.94, 0.46))
    depth4tom = crop_image("depth4tom_p1.png", "depth4tom_core.png", (0.52, 0.28, 0.94, 0.59))
    remake = crop_image("remake_p1.png", "remake_core.png", (0.49, 0.46, 0.94, 0.65))
    layered = crop_image("layereddepth_p2.png", "layereddepth_core.png", (0.08, 0.055, 0.92, 0.37))
    seegroup = crop_image("seegroup_p2.png", "seegroup_core.png", (0.08, 0.055, 0.92, 0.25))

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    blank = prs.slide_layouts[6]

    # Slide 1 — thesis
    slide = prs.slides.add_slide(blank)
    background(slide)
    textbox(slide, "透明物体抓取", 0.72, 0.72, 3.3, 0.28, 11, TEAL, True)
    textbox(slide, "从单一深度到多界面几何", 0.72, 1.35, 10.8, 0.60, 32, INK, True)
    textbox(slide, "相关工作 · Idea · 复现进度 · 下一步", 0.74, 2.12, 8.0, 0.28, 15, MUTED)
    line(slide, 0.74, 2.82, 12.56, 2.82, LINE, 1.0)

    # One-ray diagram, deliberately minimal.
    textbox(slide, "普通深度", 0.82, 3.58, 1.35, 0.22, 12, MUTED, True)
    line(slide, 2.18, 3.68, 5.15, 3.68, BLUE, 2.0)
    circle(slide, 3.10, 3.56, 0.24, BLUE)
    textbox(slide, "一个像素 = 一个 z", 2.45, 4.07, 2.45, 0.23, 14, BLUE, True, PP_ALIGN.CENTER)

    textbox(slide, "我们的目标", 6.10, 3.58, 1.45, 0.22, 12, MUTED, True)
    line(slide, 7.55, 3.68, 11.83, 3.68, TEAL, 2.0)
    for offset, label in [(0.45, "外壁"), (1.25, "内壁"), (2.55, "内壁"), (3.35, "外壁")]:
        circle(slide, 7.55 + offset, 3.55, 0.26, TEAL)
        textbox(slide, label, 7.40 + offset, 4.02, 0.56, 0.18, 8.5, MUTED, False, PP_ALIGN.CENTER)
    textbox(slide, "一个像素 = 有序界面事件", 8.05, 4.41, 3.42, 0.25, 14, TEAL, True, PP_ALIGN.CENTER)
    rect(slide, 1.05, 5.57, 11.20, 0.62, SOFT_TEAL, SOFT_TEAL, rounded=True)
    textbox(slide, "核心假设：保留外壁、内壁和空腔，能减少抓取规划中的穿壁与“实心化”错误。",
            1.24, 5.76, 10.82, 0.23, 15, INK, True, PP_ALIGN.CENTER)
    textbox(slide, "阶段汇报 · 2026.08", 0.73, 6.84, 2.5, 0.17, 8.3, MUTED)

    # Slide 2 — hardware progress
    slide = prs.slides.add_slide(blank)
    background(slide)
    title(slide, "硬件进度：从 D455 深度到 PIPER-X 基座已打通", 2)
    textbox(slide, "PIPER-X + AGILE.X 夹爪 + Intel RealSense D455 · 实机验证 2026.08.09",
            0.66, 1.13, 9.60, 0.24, 12, MUTED)
    chip(slide, "同步快照已通过", 10.82, 1.08, 1.80, SOFT_TEAL, GREEN)

    hardware_chain = [
        (0.72, "末端 D455", "RGB-D 848×480 @ 30 Hz\nUSB 3.2 · eye-in-hand", BLUE),
        (3.75, "彩色对齐深度", "camera_p\n24 / 24 角点有效", BLUE),
        (6.78, "手眼外参", "flange_T_camera\nPark · 15 + 3 姿态", TEAL),
        (9.81, "PIPER-X 基座", "base_T_flange\n法兰反馈 + FK", GREEN),
    ]
    for x, heading, body, accent in hardware_chain:
        rect(slide, x, 1.62, 2.42, 1.42, WHITE, LINE, rounded=True)
        rect(slide, x, 1.62, 0.07, 1.42, accent, accent)
        textbox(slide, heading, x + 0.22, 1.90, 1.98, 0.25,
                13, INK, True, PP_ALIGN.CENTER)
        textbox(slide, body, x + 0.20, 2.34, 2.02, 0.42,
                9.7, MUTED, False, PP_ALIGN.CENTER)
    for x in (3.25, 6.28, 9.31):
        arrow = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(x), Inches(2.04), Inches(0.32), Inches(0.48)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = color(LINE)
        arrow.line.color.rgb = color(LINE)

    rect(slide, 1.46, 3.34, 10.40, 0.58, SOFT_TEAL, SOFT_TEAL, rounded=True)
    textbox(slide, "base_p = base_T_flange @ flange_T_camera @ camera_p",
            1.69, 3.51, 9.94, 0.22, 14, TEAL, True, PP_ALIGN.CENTER)

    metrics = [
        (0.72, "统一时间戳", "5.420 ms", "相机曝光 ↔ 法兰反馈", GREEN),
        (3.84, "手眼独立留出", "3.660 mm / 0.663°", "3 个未参与求解姿态", TEAL),
        (6.96, "深度板面拟合", "0.692 mm RMS", "30 mm 边长误差 0.307 mm", BLUE),
        (10.08, "基座位姿一致性", "3.469 mm / 0.513°", "深度点 → 机器人 Base", ORANGE),
    ]
    for x, heading, value, note, accent in metrics:
        rect(slide, x, 4.28, 2.53, 1.24, WHITE, LINE, rounded=True)
        textbox(slide, heading, x + 0.18, 4.48, 2.17, 0.19,
                10.2, MUTED, True, PP_ALIGN.CENTER)
        textbox(slide, value, x + 0.14, 4.82, 2.25, 0.23,
                13.2, accent, True, PP_ALIGN.CENTER)
        textbox(slide, note, x + 0.14, 5.18, 2.25, 0.16,
                8.3, MUTED, False, PP_ALIGN.CENTER)

    rect(slide, 0.95, 5.91, 11.43, 0.71, SOFT_BLUE, SOFT_BLUE, rounded=True)
    textbox(slide, "下一步", 1.21, 6.14, 1.05, 0.20, 11.5, BLUE, True)
    textbox(slide, "连续多帧记录 + 法兰位姿插值，接透明物体 mask、目标位姿与抓取。",
            2.20, 6.10, 9.78, 0.28, 12.8, INK, True)
    footer(slide, "统一静态快照已输出 308,677 点 Base PLY；相机支架不可松动。")

    # Slide 3 — single-depth work
    slide = prs.slides.add_slide(blank)
    background(slide)
    title(slide, "相关论文①：单层深度路线", 3)
    cards = [
        (0.60, "TransCG", "RGB-D 深度补全 + 抓取数据", transcg, "RAL 2022"),
        (4.48, "Depth4ToM", "透明 / 镜面单目深度估计", depth4tom, "ICCV 2023"),
        (8.36, "ReMake", "mask 辅助的 metric depth + 抓取", remake, "2026"),
    ]
    for x, name, desc, image, year in cards:
        textbox(slide, name, x, 1.25, 2.4, 0.27, 16, INK, True)
        chip(slide, year, x + 2.51, 1.21, 0.75, SOFT_BLUE, BLUE)
        textbox(slide, desc, x, 1.63, 3.36, 0.40, 10.8, MUTED)
        picture_contain(slide, image, x, 2.16, 3.36, 2.65)
    rect(slide, 0.82, 5.41, 11.70, 0.76, SOFT, SOFT, rounded=True)
    textbox(slide, "共同点", 1.10, 5.64, 1.12, 0.20, 12, TEAL, True)
    textbox(slide, "都把输出压缩成一张深度图；对透明物体有用，但无法直接表达内壁与空腔。",
            2.18, 5.61, 9.87, 0.24, 14, INK, True)
    textbox(slide, "这也是 DepthHypothesisPack 要突破的表示瓶颈。",
            1.10, 6.43, 10.5, 0.25, 13, TEAL, True)
    footer(slide, "图片截自本地论文 PDF：TransCG、Depth4ToM、ReMake。")

    # Slide 4 — multi-layer work
    slide = prs.slides.add_slide(blank)
    background(slide)
    title(slide, "相关论文②：多层深度路线", 4)
    textbox(slide, "LayeredDepth", 0.66, 1.22, 2.4, 0.28, 16, INK, True)
    textbox(slide, "提出真实 / 合成 multi-layer benchmark，并定义 layer_first 与 layer_all 评测。",
            0.66, 1.60, 5.75, 0.42, 11, MUTED)
    picture_contain(slide, layered, 0.66, 2.12, 5.72, 3.26)

    textbox(slide, "SeeGroup", 6.91, 1.22, 2.4, 0.28, 16, INK, True)
    textbox(slide, "沿同一条 ray 预测多个深度，并让模型自适应地组织这些层。",
            6.91, 1.60, 5.75, 0.42, 11, MUTED)
    picture_contain(slide, seegroup, 6.91, 2.12, 5.72, 3.26)

    rect(slide, 1.18, 5.84, 10.98, 0.67, SOFT_TEAL, SOFT_TEAL, rounded=True)
    textbox(slide, "它们证明“多界面可以被感知”；仍未回答“多界面能否让抓取规划更可靠”。",
            1.42, 6.04, 10.50, 0.23, 14, INK, True, PP_ALIGN.CENTER)
    footer(slide, "图片截自本地论文 PDF：LayeredDepth、SeeGroup。")

    # Slide 5 — candidate representation
    slide = prs.slides.add_slide(blank)
    background(slide)
    title(slide, "候选技术路线：多界面薄壳表示", 5)
    textbox(slide, "每个像素不再只有一个深度，而是保留由近到远的界面、类型和不确定性。",
            0.61, 1.18, 10.8, 0.26, 13.5, MUTED)
    chip(slide, "候选名称：DepthHypothesisPack", 9.92, 1.13, 2.70, SOFT_TEAL, TEAL)

    steps = [
        (0.74, 1.78, 2.42, "输入", "RGB + 原始深度\n+ 物体区域", BLUE),
        (4.00, 1.78, 4.18, "多界面事件包", "[深度, 存在, 类型, 不确定性] × K", TEAL),
        (9.02, 1.78, 3.58, "薄壳几何与规划", "外壁 / 内壁 / 空腔 / 开口\n碰撞与可达性检查", GREEN),
    ]
    for x, y, w, heading, body, accent in steps:
        rect(slide, x, y, w, 1.42, WHITE, LINE, rounded=True)
        rect(slide, x, y, 0.07, 1.42, accent, accent)
        textbox(slide, heading, x + 0.24, y + 0.25, w - 0.47, 0.27,
                14, INK, True, PP_ALIGN.CENTER)
        textbox(slide, body, x + 0.24, y + 0.76, w - 0.47, 0.41,
                10.5, MUTED, False, PP_ALIGN.CENTER)
    for x in (3.34, 8.34):
        arrow = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(x), Inches(2.20), Inches(0.40), Inches(0.54)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = color(LINE)
        arrow.line.color.rgb = color(LINE)

    rect(slide, 0.88, 3.55, 11.57, 1.54, SOFT, SOFT, rounded=True)
    textbox(slide, "同一像素的一条相机射线", 1.15, 3.81, 2.10, 0.22, 11.5, INK, True)
    line(slide, 3.05, 4.18, 11.58, 4.18, MUTED, 1.3)
    line(slide, 4.12, 4.18, 5.28, 4.18, ORANGE, 5.0)
    line(slide, 5.28, 4.18, 8.90, 4.18, TEAL, 5.0)
    line(slide, 8.90, 4.18, 10.06, 4.18, ORANGE, 5.0)
    ray_events = [
        (3.99, "d¹", "空气→材料", BLUE),
        (5.15, "d²", "材料→空腔", ORANGE),
        (8.77, "d³", "空腔→材料", TEAL),
        (9.93, "d⁴", "材料→空气", GREEN),
    ]
    for x, depth, transition, accent in ray_events:
        circle(slide, x, 4.04, 0.28, accent)
        textbox(slide, depth, x - 0.10, 3.72, 0.48, 0.18,
                10, accent, True, PP_ALIGN.CENTER)
        textbox(slide, transition, x - 0.34, 4.55, 0.96, 0.18,
                8.2, MUTED, False, PP_ALIGN.CENTER)
    textbox(slide, "材料", 4.34, 4.53, 0.72, 0.18, 8.6, ORANGE, True, PP_ALIGN.CENTER)
    textbox(slide, "空腔", 6.77, 4.53, 0.72, 0.18, 8.6, TEAL, True, PP_ALIGN.CENTER)
    textbox(slide, "材料", 9.13, 4.53, 0.72, 0.18, 8.6, ORANGE, True, PP_ALIGN.CENTER)

    rect(slide, 0.93, 5.49, 11.47, 0.90, SOFT_ORANGE, SOFT_ORANGE, rounded=True)
    textbox(slide, "当前边界", 1.21, 5.78, 1.28, 0.20, 11.5, ORANGE, True)
    textbox(slide, "表示定义与评测工具已经实现；多界面预测网络尚未训练，抓取收益尚未验证。",
            2.42, 5.74, 9.55, 0.28, 13.2, INK, True)
    footer(slide, "本页是候选方法设计，不是已经完成的模型结果。")

    # Slide 6 — implemented ShellBench infrastructure
    slide = prs.slides.add_slide(blank)
    background(slide)
    title(slide, "已经实现：ShellBench 统一表示与评测工具", 6)
    textbox(slide, "它不是一个新模型，而是让单层、多层和完整真值在同一规则下公平比较。",
            0.66, 1.16, 10.9, 0.25, 13.2, MUTED)
    chip(slide, "代码与测试已完成", 10.84, 1.11, 1.80, SOFT_TEAL, GREEN)

    sources = [
        (1.64, "单层深度", "K = 1；不虚构后表面", BLUE),
        (2.55, "多层预测", "K 可变；保留深度顺序", TEAL),
        (3.46, "完整真值", "TablewareNet 射线求交", ORANGE),
    ]
    for y, heading, note, accent in sources:
        rect(slide, 0.67, y, 2.61, 0.68, WHITE, LINE, rounded=True)
        rect(slide, 0.67, y, 0.07, 0.68, accent, accent)
        textbox(slide, heading, 0.91, y + 0.13, 1.02, 0.20, 10.8, INK, True)
        textbox(slide, note, 1.91, y + 0.11, 1.11, 0.40, 7.9, MUTED, False, PP_ALIGN.RIGHT)

    arrow = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(3.45), Inches(2.55), Inches(0.38), Inches(0.54)
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = color(LINE)
    arrow.line.color.rgb = color(LINE)

    rect(slide, 3.99, 1.92, 3.45, 2.03, SOFT_TEAL, SOFT_TEAL, rounded=True)
    textbox(slide, "统一 RayEvents", 4.30, 2.19, 2.83, 0.25, 14.2, TEAL, True, PP_ALIGN.CENTER)
    textbox(slide, "depths_m        [K,H,W]\nvalid_mask       [K,H,W]\ntransition_type  [K,H,W]\nuncertainty_m    可选",
            4.42, 2.67, 2.59, 0.89, 10.2, INK, False, PP_ALIGN.LEFT)

    arrow = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(7.60), Inches(2.55), Inches(0.38), Inches(0.54)
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = color(LINE)
    arrow.line.color.rgb = color(LINE)

    rect(slide, 8.14, 1.64, 4.53, 2.70, WHITE, LINE, rounded=True)
    textbox(slide, "统一评分", 8.48, 1.93, 3.85, 0.25, 14.2, INK, True, PP_ALIGN.CENTER)
    metric_rows = [
        ("界面是否找全", "Precision / Recall / F1"),
        ("界面数量是否正确", "Event-count accuracy"),
        ("深度偏差", "匹配后的 MAE / RMSE"),
        ("材料转换是否正确", "Transition / topology"),
    ]
    for i, (meaning, metric) in enumerate(metric_rows):
        y = 2.40 + i * 0.43
        textbox(slide, meaning, 8.48, y, 1.82, 0.18, 9.4, MUTED)
        textbox(slide, metric, 10.22, y, 2.02, 0.18, 9.4, BLUE, True, PP_ALIGN.RIGHT)

    rect(slide, 0.93, 4.73, 11.47, 0.60, SOFT, SOFT, rounded=True)
    textbox(slide, "合法拓扑：空气 → 材料 → 空腔 → 材料 → 空气",
            1.22, 4.93, 4.69, 0.21, 11.6, TEAL, True)
    textbox(slide, "坐标链：(u,v,dₖ) → 相机点 → 手眼外参 → PIPER Base 薄壳",
            6.02, 4.93, 5.99, 0.21, 10.8, INK, True, PP_ALIGN.RIGHT)

    rect(slide, 0.93, 5.65, 5.48, 0.76, SOFT_TEAL, SOFT_TEAL, rounded=True)
    textbox(slide, "已经完成", 1.20, 5.88, 1.10, 0.20, 11.4, GREEN, True)
    textbox(slide, "数据格式 · 转换器 · 真值生成 · 评测 · 单元/集成测试",
            2.23, 5.86, 3.86, 0.27, 9.8, INK, True)
    rect(slide, 6.62, 5.65, 5.78, 0.76, SOFT_ORANGE, SOFT_ORANGE, rounded=True)
    textbox(slide, "尚未完成", 6.90, 5.88, 1.10, 0.20, 11.4, ORANGE, True)
    textbox(slide, "预测网络 · 三维薄壳提升 · 规划器对照 · 实机闭环",
            7.93, 5.86, 4.13, 0.27, 9.8, INK, True)
    footer(slide, "ShellBench 已能评价表示质量，但尚未产生机器人抓取成功率结论。")

    # Slide 7 — reproduction status in plain language
    slide = prs.slides.add_slide(blank)
    background(slide)
    title(slide, "复现进度：已经证明什么，还缺什么", 7)

    textbox(slide, "已完成的证据", 0.66, 1.27, 2.3, 0.24, 14, INK, True)
    chip(slide, "结果可信", 6.35, 1.22, 1.22, SOFT_TEAL, GREEN)
    completed = [
        ("单层深度基线", "Booster · 228 张", "均方根误差 136.28 mm（约 13.6 cm）", "与论文表格一致", GREEN),
        ("多层深度对比", "LayeredDepth · 300 张", "4 个界面顺序全对：29.95% → 72.41%", "多层模型明显更完整", TEAL),
        ("薄壳几何检查", "TablewareNet · 100 个场景", "已知真实物体区域时，界面识别 F1 = 79.15%", "说明几何路线可行，不是论文正式分数", ORANGE),
    ]
    for i, (heading, scope, result, note, accent) in enumerate(completed):
        y = 1.75 + i * 1.23
        rect(slide, 0.66, y, 7.14, 1.08, WHITE, LINE, rounded=True)
        rect(slide, 0.66, y, 0.07, 1.08, accent, accent)
        textbox(slide, heading, 0.92, y + 0.17, 2.08, 0.21, 11.5, INK, True)
        textbox(slide, scope, 0.92, y + 0.56, 2.08, 0.18, 9.0, MUTED)
        textbox(slide, result, 3.05, y + 0.16, 4.30, 0.23, 11.1, accent, True, PP_ALIGN.RIGHT)
        textbox(slide, note, 3.05, y + 0.57, 4.30, 0.18, 8.9, MUTED, False, PP_ALIGN.RIGHT)

    textbox(slide, "仍未完成", 8.06, 1.27, 2.1, 0.24, 14, INK, True)
    chip(slide, "暂不下结论", 11.27, 1.22, 1.40, SOFT_ORANGE, ORANGE)
    pending = [
        ("TransCG 全量测试", "需要 52 个场景 / 23,524 个样本", "已取得 40 / 52 个场景", "缺少 12 个场景，暂不报最终指标"),
        ("机器人抓取收益", "还没有公平比较单层与多层输入", "尚无抓取成功率结论", "下一页用同一规划器直接验证"),
    ]
    for i, (heading, scope, result, note) in enumerate(pending):
        y = 1.75 + i * 1.91
        rect(slide, 8.06, y, 4.61, 1.70, WHITE, LINE, rounded=True)
        rect(slide, 8.06, y, 0.07, 1.70, ORANGE, ORANGE)
        textbox(slide, heading, 8.34, y + 0.21, 4.01, 0.23, 12.2, INK, True)
        textbox(slide, scope, 8.34, y + 0.60, 4.01, 0.20, 9.3, MUTED)
        textbox(slide, result, 8.34, y + 0.98, 4.01, 0.22, 11.3, ORANGE, True)
        textbox(slide, note, 8.34, y + 1.34, 4.01, 0.18, 8.9, MUTED)

    rect(slide, 0.86, 5.72, 11.58, 0.86, SOFT_BLUE, SOFT_BLUE, rounded=True)
    textbox(slide, "目前结论", 1.14, 6.00, 1.38, 0.22, 12, BLUE, True)
    textbox(slide, "多层深度能更完整地描述透明物体；是否让机器人抓得更稳，还需要下一步实验。",
            2.41, 5.96, 9.57, 0.31, 13.2, INK, True)
    footer(slide, "所有数字都来自本地完整运行；未完成的实验不提前写成结论。")

    # Slide 8 — next work in plain language
    slide = prs.slides.add_slide(blank)
    background(slide)
    title(slide, "下一步：证明“多层深度是否真的更好抓”", 8)
    textbox(slide, "保持物体、物体区域和抓取规划器完全相同，只改变输入的几何信息。",
            0.66, 1.17, 11.4, 0.25, 13.2, MUTED)

    rect(slide, 1.05, 1.66, 11.23, 0.69, SOFT_TEAL, SOFT_TEAL, rounded=True)
    textbox(slide, "核心问题：只知道最前表面，与知道前后 / 内外表面相比，哪一种抓取错误更少？",
            1.30, 1.88, 10.73, 0.25, 14.2, TEAL, True, PP_ALIGN.CENTER)

    stages = [
        (1, "做公平对比", "方案 A：只给最前表面\n方案 B：给出完整多层界面", "其余条件全部相同", BLUE),
        (2, "统计三类错误", "① 穿过物体外壁\n② 把空腔当成实心\n③ 抓取位姿碰撞或不可达", "重复实验，比较错误率和波动", ORANGE),
        (3, "根据结果决定", "错误明显减少：训练我们的多层模型\n没有减少：先修改规划器或评测", "不盲目堆模型", GREEN),
    ]
    for i, (number, heading, body, note, accent) in enumerate(stages):
        x = 0.74 + i * 4.18
        rect(slide, x, 2.77, 3.77, 2.27, WHITE, LINE, rounded=True)
        circle(slide, x + 0.22, 3.02, 0.44, accent, str(number))
        textbox(slide, heading, x + 0.82, 3.06, 2.62, 0.25, 14, INK, True)
        textbox(slide, body, x + 0.30, 3.60, 3.17, 0.75,
                10.7, MUTED, False, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        rect(slide, x + 0.31, 4.50, 3.15, 0.34, SOFT, SOFT, rounded=True)
        textbox(slide, note, x + 0.39, 4.58, 2.99, 0.18,
                8.7, accent, True, PP_ALIGN.CENTER)

    rect(slide, 0.96, 5.43, 11.41, 0.64, SOFT_BLUE, SOFT_BLUE, rounded=True)
    textbox(slide, "同时补齐已有方法：TransCG 剩余 12 个场景，并完成 DFNet / ReMake 全量测试。",
            1.25, 5.63, 10.83, 0.23, 12.0, BLUE, True, PP_ALIGN.CENTER)
    textbox(slide, "最终要证明的不是“深度图更漂亮”，而是“机械臂少犯抓取错误”。",
            1.74, 6.46, 9.86, 0.28, 15.4, TEAL, True, PP_ALIGN.CENTER)
    footer(slide, "先用真实几何验证价值，再投入训练多界面预测模型。")

    # Slide 9 — technical appendix
    slide = prs.slides.add_slide(blank)
    background(slide)
    title(slide, "技术备用页：匹配、指标与因果对照", 9)
    chip(slide, "追问时展示", 11.18, 1.12, 1.46, SOFT_BLUE, BLUE)

    rect(slide, 0.66, 1.48, 5.86, 2.94, WHITE, LINE, rounded=True)
    textbox(slide, "1｜有序界面匹配", 0.96, 1.78, 2.75, 0.25, 14, INK, True)
    textbox(slide, "真值", 0.98, 2.30, 0.55, 0.20, 9.5, MUTED, True)
    textbox(slide, "0.400   0.405   0.490   0.495 m", 1.63, 2.28, 3.95, 0.22,
            11.0, TEAL, True, PP_ALIGN.CENTER)
    textbox(slide, "预测", 0.98, 2.76, 0.55, 0.20, 9.5, MUTED, True)
    textbox(slide, "0.402      缺失      0.492   0.497 m", 1.63, 2.74, 3.95, 0.22,
            11.0, ORANGE, True, PP_ALIGN.CENTER)
    line(slide, 1.80, 3.18, 5.49, 3.18, LINE, 0.9)
    textbox(slide, "动态规划按三个规则匹配：", 0.98, 3.42, 2.42, 0.20, 10.2, INK, True)
    textbox(slide, "① 保持由近到远　② 先让匹配数最多　③ 再让总误差最小",
            0.98, 3.80, 5.09, 0.22, 10.0, MUTED)
    textbox(slide, "只有误差 ≤ δ 的界面可匹配；δ 可配置，示例为 5 mm。",
            0.98, 4.12, 5.09, 0.18, 8.8, BLUE, True)

    rect(slide, 6.81, 1.48, 5.86, 2.94, WHITE, LINE, rounded=True)
    textbox(slide, "2｜评分指标", 7.11, 1.78, 2.40, 0.25, 14, INK, True)
    rect(slide, 7.10, 2.27, 2.46, 1.39, SOFT_BLUE, SOFT_BLUE, rounded=True)
    textbox(slide, "P = 匹配数 / 预测数\nR = 匹配数 / 真值数\nF1 = 2PR / (P + R)",
            7.34, 2.55, 1.98, 0.72, 10.4, BLUE, True, PP_ALIGN.CENTER)
    rect(slide, 9.82, 2.27, 2.46, 1.39, SOFT_TEAL, SOFT_TEAL, rounded=True)
    textbox(slide, "界面数量正确率\n匹配深度 MAE / RMSE\n类型 F1 / 拓扑合法率",
            10.05, 2.53, 2.00, 0.78, 10.0, TEAL, True, PP_ALIGN.CENTER)
    textbox(slide, "漏掉后层会降低 Recall；多报假界面会降低 Precision，不能只看匹配后的深度误差。",
            7.16, 3.90, 5.14, 0.27, 9.3, MUTED, True, PP_ALIGN.CENTER)

    textbox(slide, "3｜三组因果对照", 0.70, 4.73, 2.84, 0.25, 14, INK, True)
    comparisons = [
        (0.70, "A｜GT-front + 固定规划器", "只给最前表面\n建立单层基准", BLUE),
        (4.50, "B｜GT-full + 同一规划器", "只替换完整薄壳几何\n隔离“表示”的收益", TEAL),
        (8.30, "C｜GT-full + 薄壳规划器", "表示与规划共同正确\n估计系统上界", GREEN),
    ]
    for x, heading, body, accent in comparisons:
        rect(slide, x, 5.14, 3.56, 1.24, WHITE, LINE, rounded=True)
        rect(slide, x, 5.14, 0.07, 1.24, accent, accent)
        textbox(slide, heading, x + 0.22, 5.39, 3.10, 0.23,
                11.1, INK, True, PP_ALIGN.CENTER)
        textbox(slide, body, x + 0.26, 5.83, 3.02, 0.39,
                9.4, MUTED, False, PP_ALIGN.CENTER)
    footer(slide, "对象、候选抓取、规划器参数和数据划分全部冻结；GT 只用于离线审计，不参与候选生成。")

    prs.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    build()
